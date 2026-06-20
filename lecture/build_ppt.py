#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
전반부 AI원론 PPT 생성기 — build_ppt.py (v2.1)
================================================
3인 비판 검토 반영 **순수 python-pptx 드로잉**(gpt-image 의존 최소).
v2.1: 누락 템플릿 5종(T3/T5_cards/T5_chips/T8_harness/T8_security) 추가 + stub 7장 완전 구현.
       → 17장 전부 템플릿 기반 완성(더 이상 TODO stub 없음).

실행:  python build_ppt.py
산출:  전반부_AI원론_PPT.pptx (17장, 16:9)
의존:  pip install python-pptx Pillow

설계 원칙 (검토 반영):
 - 도형 전부 python-pptx 코드로 그림 → 도형/텍스트 같은 좌표계 → 정합 100%, 결정적
 - 한글 자수 공식 계수 0.52(한글 음절 advance ≈1em = 영문 2배). chars_per_line() 참조.
 - 색 위계: ACCENT(주강조·해법) + WARN(환각·별도검토) — 단색 정책 + 명도 위계
 - 강의장 가독: 본문 20pt, 라벨 13pt, 푸터 11pt
 - 정직성: S07/S08 의도적 왜곡 예시 표기, S09 '확률적으로 찍는다' 금지 부표기, S14 감사 본류 미구축
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ============================================================
# A. 전역 디자인 토큰
# ============================================================
BG          = RGBColor(0xFF, 0xFF, 0xFF)   # 순백 배경
INK         = RGBColor(0x1A, 0x1A, 0x1A)   # 제목/검정
BODY        = RGBColor(0x4A, 0x4A, 0x4A)   # 본문 회색
MUTE        = RGBColor(0x6B, 0x72, 0x80)   # 보조/출처
LINE_C      = RGBColor(0xE5, 0xE5, 0xE5)   # 구분선/표 테두리
ACCENT      = RGBColor(0xFF, 0x4D, 0x30)   # 주강조(밝은 오렌지) — 해법·신뢰·주 포인트
WARN        = RGBColor(0xC2, 0x41, 0x0C)   # 경고(어두운 오렌지) — 환각·별도검토 (명도 위계)
WARN_SOFT   = RGBColor(0xFB, 0xEE, 0xE9)   # 경고 채움
ACCENT_SOFT = RGBColor(0xFF, 0xF1, 0xED)   # 오렌지 채움
PANEL_BG    = RGBColor(0xFA, 0xFA, 0xFA)   # placeholder/패널 배경

FONT  = "Pretendard"        # 한글+영문 통일. 없으면 시스템 fallback.
SW, SH = 13.333, 7.5        # 16:9 inches

PT_DISPLAY, PT_H1, PT_H2 = 56, 30, 22
PT_BODY, PT_BODY_SM, PT_CAPTION = 20, 16, 14
PT_LABEL, PT_FOOTER = 13, 11

# ============================================================
# 헬퍼
# ============================================================
def pct(v, dim):
    return Inches(v / 100.0 * dim)

def chars_per_line(box_w_pct, pt):
    """1줄 안전 한글 자수(검토 반영: 계수 0.52). 초과 시 박스 확장 또는 pt 축소."""
    w_in = box_w_pct / 100.0 * SW
    return int(w_in * 72 / pt * 0.52)

def S(text, size=PT_BODY, bold=False, color=BODY, italic=False, font=FONT):
    return (text, {"size": size, "bold": bold, "color": color, "italic": italic, "font": font})

def _apply_run(r, text, st):
    r.text = text
    r.font.size = Pt(st.get("size", PT_BODY))
    r.font.bold = st.get("bold", False)
    r.font.italic = st.get("italic", False)
    r.font.name = st.get("font", FONT)
    r.font.color.rgb = st.get("color", BODY)

def add_box(slide, x, y, w, h, lines,
            shape=None, fill=None, line=None, line_w=1.0,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            line_spacing=None, wrap=True, margin=0.06):
    """도형(optional)+텍스트 통합 박스. 도형을 쓰면 도형-텍스트 같은 좌표(정합 100%).
    lines: List[paragraph]. paragraph: List[run] 또는 단일 run. run: (text, style)."""
    left, top, width, height = pct(x, SW), pct(y, SH), pct(w, SW), pct(h, SH)
    if shape is not None:
        sp = slide.shapes.add_shape(shape, left, top, width, height)
        if fill is None: sp.fill.background()
        else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None: sp.line.fill.background()
        else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
        try: sp.shadow.inherit = False
        except Exception: pass
    else:
        sp = slide.shapes.add_textbox(left, top, width, height)
    tf = sp.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin * 0.6); tf.margin_bottom = Inches(margin * 0.6)
    tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        if isinstance(para, tuple): para = [para]
        for run in para:
            _apply_run(p.add_run(), run[0], run[1])
    return sp

def add_line(slide, x1, y1, x2, y2, color=LINE_C, w=1.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    pct(x1, SW), pct(y1, SH), pct(x2, SW), pct(y2, SH))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    if dash:
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash', {'val': dash}))
    return cn

def add_arrow(slide, x1, y1, x2, y2, color=ACCENT, w=2.0):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    pct(x1, SW), pct(y1, SH), pct(x2, SW), pct(y2, SH))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd',
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn

def draw_chrome(slide, page_no, section_kr, label_en, source=None):
    """본문 슬라이드 공통 크롬(Nexus 시그니처): 로고·블록라벨·구분선·페이지번호·출처."""
    add_box(slide, 6, 4.5, 14, 4, [[S("원론", PT_LABEL, True, INK), S(".", PT_LABEL, True, ACCENT)]])
    add_box(slide, 58, 4.5, 36, 3.5, [[S(label_en, PT_LABEL, True, ACCENT)]], align=PP_ALIGN.RIGHT)
    add_line(slide, 6, 10, 94, 10, LINE_C, 1.0)
    add_box(slide, 6, 93, 50, 4, [[
        S(f"{page_no:02d} / 17", PT_FOOTER, False, MUTE),
        S(f"   ·   {section_kr}", PT_FOOTER, False, MUTE),
    ]])
    if source:
        add_box(slide, 46, 93, 48, 4, [[S(source, PT_FOOTER, False, MUTE)]], align=PP_ALIGN.RIGHT)

def blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s

# ============================================================
# 템플릿 도형 함수 (T1~T9)
# ============================================================

def T1_cover(slide):
    """S01 표지/지도 — 두 개의 30분."""
    add_box(slide, 6, 10, 30, 6, [[S("원론", PT_DISPLAY, True, INK), S(".", PT_DISPLAY, True, ACCENT)]])
    add_box(slide, 70, 12, 24, 3, [[S("회계사 AI 원론 · 2026", PT_LABEL, False, MUTE)]], align=PP_ALIGN.RIGHT)
    add_line(slide, 6, 20, 94, 20, LINE_C, 1.0)
    add_box(slide, 6, 28, 88, 7, [[S("두 개의 30분으로 봅니다", PT_H1, True, INK)]])
    add_box(slide, 6, 44, 41, 30,
            [[S("전반부 30분", PT_LABEL, True, BODY)], [S("개념 PPT", PT_H2, True, INK)],
             [S("왜 틀리는지 · 어떻게 쓰는지", PT_CAPTION, False, MUTE)]],
            shape=MSO_SHAPE.RECTANGLE, line=LINE_C, line_w=1.0, anchor=MSO_ANCHOR.MIDDLE, margin=0.18)
    add_arrow(slide, 47, 59, 53, 59, ACCENT, 2.0)
    add_box(slide, 53, 44, 41, 30,
            [[S("후반부 30분", PT_LABEL, True, ACCENT)], [S("Claude Code 라이브", PT_H2, True, INK)],
             [S("직접 도구를 켜서 시연", PT_CAPTION, False, MUTE)]],
            shape=MSO_SHAPE.RECTANGLE, line=ACCENT, line_w=1.5, anchor=MSO_ANCHOR.MIDDLE, margin=0.18)
    add_box(slide, 6, 80, 88, 5,
            [[S("그리고 뒤이어 1시간 — 완성된 회계 감사 시스템 실연", PT_CAPTION, False, MUTE)]])
    add_box(slide, 6, 93, 30, 4, [[S("00 / 17   ·   표지", PT_FOOTER, False, MUTE)]])

def T2_statement(slide, label_en, big_lines, sub=None):
    """S02/S03 명제·빅스테이트먼트형."""
    add_box(slide, 6, 16, 70, 3.5, [[S(label_en, PT_LABEL, True, ACCENT)]])
    add_box(slide, 6, 28, 88, 44, big_lines, line_spacing=1.35, align=PP_ALIGN.LEFT)
    if sub:
        add_box(slide, 6, 86, 88, 4, [[S(sub, PT_CAPTION, False, MUTE)]])

def T3_timeline(slide, h1, stages, active=None, dashed=None, caption=None):
    """S04 가로 타임라인. stages: [(num, label, desc)]. active=ACCENT 노드, dashed=점선 노드."""
    add_box(slide, 6, 14, 80, 5, [[S(h1, PT_H1, True, INK)]])
    n = len(stages)
    gap = 78 / n
    centers = [11 + gap * (i + 0.5) for i in range(n)]
    cy = 42
    add_line(slide, centers[0], cy, centers[-1], cy, LINE_C, 1.5)
    if active is not None and active < n - 1:
        add_line(slide, centers[active], cy, centers[active + 1], cy, ACCENT, 2.5)
    if dashed is not None and dashed < n - 1:
        add_line(slide, centers[dashed - 1] if dashed > 0 else centers[0], cy, centers[dashed], cy, MUTE, 1.5, dash="dash")
    for i, (num, label, desc) in enumerate(stages):
        cx = centers[i]
        is_active = (i == active)
        col = ACCENT if is_active else MUTE
        fill = ACCENT if is_active else BG
        add_box(slide, cx - 3, cy - 3, 6, 6, [[S(num, 16, True, INK if is_active else MUTE)]],
                shape=MSO_SHAPE.OVAL, fill=fill, line=col, line_w=1.5,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_box(slide, cx - 10, cy + 5, 20, 3.5, [[S(label, PT_LABEL, True, col)]], align=PP_ALIGN.CENTER)
        desc_list = desc if isinstance(desc, list) else [desc]
        add_box(slide, cx - 10, cy + 10, 20, 14,
                [[S(d, 11, False, BODY)] for d in desc_list],
                align=PP_ALIGN.CENTER, line_spacing=1.2)
    if caption:
        add_box(slide, 6, 82, 88, 5, [[S(caption, PT_CAPTION, True, MUTE)]], align=PP_ALIGN.CENTER)

def T4_concept_map(slide):
    """★ S05 개념지도 — 강의 시각적 중심. 도형+텍스트 통합(정합 100%)."""
    add_box(slide, 6, 7, 88, 5,
            [[S("AI는 언어는 잘하지만, 말로만 숫자·사실을 꺼내면 틀린다 → 숫자는 계산기·사실은 검색·판단·서명은 회계사",
                PT_CAPTION, False, INK)]],
            shape=MSO_SHAPE.RECTANGLE, fill=ACCENT_SOFT,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.12)
    cx, cy = 50, 45
    for nx, ny, col in [(19, 31, LINE_C), (19, 64, LINE_C), (81, 31, WARN), (50, 79, INK), (81, 64, LINE_C)]:
        add_line(slide, cx, cy, nx, ny, col, 1.0)
    add_box(slide, 41, 36, 18, 18,
            [[S("LLM", PT_H2, True, ACCENT)], [S("언어는 잘하는 두뇌", 11, False, BODY)]],
            shape=MSO_SHAPE.OVAL, line=ACCENT, line_w=1.5,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_box(slide, 8, 24, 22, 14,
            [[S("RAG", PT_LABEL, True, ACCENT)], [S("사실은 검색에", PT_CAPTION, False, BODY)],
             [S("조문·예규·판례", 10, False, MUTE)]],
            shape=MSO_SHAPE.RECTANGLE, line=LINE_C, line_w=1.0, anchor=MSO_ANCHOR.MIDDLE, margin=0.1)
    add_box(slide, 8, 57, 22, 14,
            [[S("코드·계산기", PT_LABEL, True, ACCENT)], [S("숫자는 계산기에", PT_CAPTION, False, BODY)],
             [S("환각 안 터짐", 10, False, MUTE)]],
            shape=MSO_SHAPE.RECTANGLE, line=LINE_C, line_w=1.0, anchor=MSO_ANCHOR.MIDDLE, margin=0.1)
    add_box(slide, 70, 24, 22, 14,
            [[S("환각", PT_LABEL, True, WARN)], [S("말로만 틀린다", PT_CAPTION, False, WARN)],
             [S("조용한 실수", 10, False, MUTE)]],
            shape=MSO_SHAPE.RECTANGLE, line=WARN, line_w=1.5, fill=WARN_SOFT,
            anchor=MSO_ANCHOR.MIDDLE, margin=0.1)
    add_box(slide, 33, 73, 34, 12,
            [[S("하네스 — 실행틀 + 고삐 2개", PT_LABEL, True, INK)],
             [S("① CLAUDE.md   ② Skill   (뒤 30분 제작)", 11, False, BODY)]],
            shape=MSO_SHAPE.RECTANGLE, line=INK, line_w=1.2, anchor=MSO_ANCHOR.MIDDLE, margin=0.1)
    add_box(slide, 70, 57, 22, 14,
            [[S("에이전트·AX", PT_LABEL, True, ACCENT)], [S("끝까지 맡긴다", PT_CAPTION, False, BODY)]],
            shape=MSO_SHAPE.RECTANGLE, line=LINE_C, line_w=1.0, anchor=MSO_ANCHOR.MIDDLE, margin=0.1)
    add_box(slide, 6, 88, 88, 4,
            [[S("비유 — 회계사무소 보조 인력: 말 잘하는 비서(LLM) · 자료 펼쳐주는 손(RAG) · 결정적 계산기(코드) · 근태·결재 위 업무지시서·절차매뉴얼(하네스) · 심부름꾼(에이전트)",
                9, False, MUTE)]], align=PP_ALIGN.CENTER)

def T5_cards(slide, h1, cards, caption=None, num_color=ACCENT):
    """S06 N-카드 가로. cards: [(num_or_None, title, desc_lines)]."""
    add_box(slide, 6, 14, 80, 5, [[S(h1, PT_H1, True, INK)]])
    n = len(cards)
    cw = (88 - 2 * (n - 1)) / n
    for i, (num, title, desc) in enumerate(cards):
        cx = 6 + i * (cw + 2)
        add_box(slide, cx, 30, cw, 50, [],
                shape=MSO_SHAPE.RECTANGLE, line=LINE_C, line_w=1.0, margin=0.18)
        lines = []
        if num: lines.append([S(num, 24, True, num_color)])
        lines.append([S(title, PT_H2, True, INK)])
        for d in (desc if isinstance(desc, list) else [desc]):
            lines.append([S(d, PT_BODY_SM, False, BODY)])
        add_box(slide, cx, 33, cw - 1, 44, lines, line_spacing=1.3, margin=0.18)
    if caption:
        add_box(slide, 6, 86, 88, 4, [[S(caption, PT_CAPTION, True, MUTE)]], align=PP_ALIGN.CENTER)

def _chip_row(slide, chips, y, fill, border):
    n = len(chips)
    cw = (88 - 2 * (n - 1)) / n
    for i, c in enumerate(chips):
        cx = 6 + i * (cw + 2)
        add_box(slide, cx, y, cw, 12, [[S(c, PT_BODY_SM, True, border)]],
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, fill=fill, line=border, line_w=1.0,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.1)

def T5_chips(slide, h1, top_chips, bottom_chips, caption=None):
    """S10 프롬프트 — 상단 4칩 + 하단 3칩."""
    add_box(slide, 6, 14, 80, 5, [[S(h1, PT_H1, True, INK)]])
    add_box(slide, 6, 24, 88, 3.5, [[S("기본 4요소", PT_LABEL, True, MUTE)]])
    _chip_row(slide, top_chips, 29, ACCENT_SOFT, ACCENT)
    add_box(slide, 6, 52, 88, 3.5, [[S("강화 3가지 (이 중 하나만 붙여도 결과가 다릅니다)", PT_LABEL, True, WARN)]])
    _chip_row(slide, bottom_chips, 57, WARN_SOFT, WARN)
    if caption:
        add_box(slide, 6, 86, 88, 4, [[S(caption, PT_CAPTION, True, ACCENT)]], align=PP_ALIGN.CENTER)

def T6_capture(slide, png_path=None, h1=None, caption_lines=None, note=None, warn=False):
    """S07/S08/S14 캡처+캡션형. png_path 없으면 placeholder 사각."""
    if h1:
        add_box(slide, 6, 13, 80, 4, [[S(h1, PT_H2, True, INK)]])
    border = WARN if warn else ACCENT
    fx, fy, fw, fh = 6, 22, 56, 62
    if png_path and os.path.exists(png_path):
        slide.shapes.add_picture(png_path, pct(fx, SW), pct(fy, SH), pct(fw, SW), pct(fh, SH))
    add_box(slide, fx, fy, fw, fh, [],
            shape=MSO_SHAPE.RECTANGLE, line=border, line_w=1.5)
    if not (png_path and os.path.exists(png_path)):
        add_box(slide, fx, fy, fw, fh,
                [[S("〔 캡처 이미지 삽입 예정 〕", PT_CAPTION, False, MUTE)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if caption_lines:
        add_box(slide, 64, 24, 30, 56, caption_lines, anchor=MSO_ANCHOR.TOP, line_spacing=1.3)
    if note:
        fill = WARN_SOFT if warn else None
        add_box(slide, 6, 86, 88, 5, [[S(note, PT_CAPTION, True, WARN if warn else MUTE)]],
                shape=MSO_SHAPE.RECTANGLE if fill else None, fill=fill, margin=0.1)

def T7_split(slide, h1, left_label, left_lines, right_label, right_lines,
             bottom=None, arrow=False, right_strong=False):
    """S09/S11/S13/S16 대비 2컬럼형. right_strong=True 시 우컬럼 INK 굵게(회계사 자리)."""
    add_box(slide, 6, 13, 80, 5, [[S(h1, PT_H1, True, INK)]])
    add_box(slide, 6, 23, 42, 3.5, [[S(left_label, PT_LABEL, True, ACCENT)]])
    add_box(slide, 6, 29, 42, 54, left_lines, line_spacing=1.3)
    add_box(slide, 52, 23, 42, 3.5, [[S(right_label, PT_LABEL, True, ACCENT)]])
    add_box(slide, 52, 29, 42, 54, right_lines, line_spacing=1.3)
    add_line(slide, 49.5, 24, 49.5, 84, LINE_C, 1.0, dash="dash")
    if arrow:
        add_arrow(slide, 47, 50, 53, 50, ACCENT, 2.0)
    if bottom:
        add_box(slide, 6, 86, 88, 5, [[S(bottom, PT_CAPTION, True, ACCENT)]],
                shape=MSO_SHAPE.RECTANGLE, fill=ACCENT_SOFT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.1)

def T8_harness(slide, h1, outer_label, inner_boxes, bottom, side_note=None):
    """S12 하네스 — 큰 사각(outer) 안 inner 고삐 박스들 + side note."""
    add_box(slide, 6, 14, 80, 5, [[S(h1, PT_H1, True, INK)]])
    add_box(slide, 8, 30, 84, 50, [[S(outer_label, PT_LABEL, True, MUTE)]],
            shape=MSO_SHAPE.RECTANGLE, line=LINE_C, line_w=1.0,
            anchor=MSO_ANCHOR.TOP, margin=0.15)
    n = len(inner_boxes)
    iw = (76 - 4 * (n - 1)) / n
    for i, (label, desc) in enumerate(inner_boxes):
        ix = 12 + i * (iw + 4)
        add_box(slide, ix, 42, iw, 30,
                [[S(f"고삐 {i + 1}", 11, True, ACCENT)], [S(label, PT_H2, True, INK)],
                 [S(desc, PT_CAPTION, False, BODY)]],
                shape=MSO_SHAPE.RECTANGLE, line=ACCENT, line_w=1.5,
                anchor=MSO_ANCHOR.MIDDLE, margin=0.15, fill=ACCENT_SOFT)
    if side_note:
        add_box(slide, 12, 74, 76, 3.5, [[S(side_note, PT_CAPTION, False, MUTE)]], align=PP_ALIGN.CENTER)
    add_box(slide, 6, 86, 88, 4, [[S(bottom, PT_CAPTION, True, INK)]], align=PP_ALIGN.CENTER)

def T8_security(slide, h1, base_lines, risks, side_note=None):
    """S15 보안 — 상단 뼈대(ACCENT_SOFT) + 하단 별도검토 3종(WARN). risks: [(num, text)]."""
    add_box(slide, 6, 14, 80, 5, [[S(h1, PT_H1, True, INK)]])
    if side_note:
        add_box(slide, 6, 21, 88, 3, [[S(side_note, PT_CAPTION, False, MUTE)]])
    add_box(slide, 6, 26, 88, 22, base_lines,
            shape=MSO_SHAPE.RECTANGLE, fill=ACCENT_SOFT, line=ACCENT, line_w=1.0,
            anchor=MSO_ANCHOR.MIDDLE, margin=0.2, line_spacing=1.3)
    add_box(slide, 6, 52, 88, 3.5, [[S("피해갈 수 없는 3종 — 별도 정보보안·컴플라이언스 검토 필요", PT_LABEL, True, WARN)]])
    cw = (88 - 2 * 2) / 3
    for i, (num, text) in enumerate(risks):
        cx = 6 + i * (cw + 2)
        add_box(slide, cx, 57, cw, 25,
                [[S(num, 20, True, WARN)], [S(text, PT_BODY_SM, False, BODY)]],
                shape=MSO_SHAPE.RECTANGLE, line=WARN, line_w=1.0, fill=WARN_SOFT,
                anchor=MSO_ANCHOR.MIDDLE, margin=0.15, line_spacing=1.25)

def T9_table(slide, h1, header, rows, highlight_row=None, bottom=None):
    """S17 진입표형(텍스트박스 그리드)."""
    add_box(slide, 6, 13, 80, 5, [[S(h1, PT_H1, True, INK)]])
    cols = [26, 14, 60]
    x0, y0, rh = 6, 26, 9
    xs = [x0, x0 + cols[0], x0 + cols[0] + cols[1]]
    for i, htext in enumerate(header):
        add_box(slide, xs[i], y0, cols[i] - 1, rh, [[S(htext, PT_LABEL, True, MUTE)]],
                shape=MSO_SHAPE.RECTANGLE, fill=RGBColor(0xF5, 0xF5, 0xF5), line=LINE_C, line_w=0.75,
                anchor=MSO_ANCHOR.MIDDLE, margin=0.1)
    for r, row in enumerate(rows):
        yy = y0 + rh * (r + 1)
        is_hi = (highlight_row is not None and r == highlight_row)
        fill = ACCENT_SOFT if is_hi else BG
        for i, cell in enumerate(row):
            runs = cell if isinstance(cell, list) else [[S(cell, PT_BODY_SM, is_hi, INK if is_hi else BODY)]]
            add_box(slide, xs[i], yy, cols[i] - 1, rh, runs,
                    shape=MSO_SHAPE.RECTANGLE, fill=fill, line=LINE_C, line_w=0.75,
                    anchor=MSO_ANCHOR.MIDDLE, margin=0.1)
    if bottom:
        add_box(slide, 6, 88, 88, 4, [[S(bottom, PT_CAPTION, True, INK)]])

# ============================================================
# 17장 조립
# ============================================================
def build(out="전반부_AI원론_PPT.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    # S01 표지
    T1_cover(blank_slide(prs))

    # S02 학습목표 + 통계
    s = blank_slide(prs)
    T2_statement(s, "LEARNING OBJECTIVE · STAT",
                 [[S("오늘이 끝나면, 업무지시서 한 장(", PT_H2, False, INK),
                   S("CLAUDE.md", PT_H2, True, ACCENT),
                   S(")을 직접 써서 AI가 내 기준·내 말투로 답하게 만들어 볼 수 있습니다.", PT_H2, False, INK)],
                  [S("약 46%", 40, True, ACCENT),
                   S("   의 회계사가 매일 AI 사용", PT_H2, False, BODY)]],
                 sub="출처: FindSkill.ai 2026 (업계 설문 기준, 공식 통계 아님)")
    draw_chrome(s, 2, "오프닝", "OBJECTIVE · STAT")

    # S03 핵심 명제 (원문 복원)
    s = blank_slide(prs)
    T2_statement(s, "THE THESIS",
                 [[S("AI는 언어는 잘하지만, ", PT_H2, False, INK),
                   S("말로만", PT_H2, True, ACCENT),
                   S(" 숫자와 사실을 꺼내면 틀린다.", PT_H2, False, INK)],
                  [S("그래서 숫자는 계산기(코드)에, 사실은 검색(RAG)에 맡기고,", PT_H2, False, INK)],
                  [S("최종 판단·서명은 회계사가 합니다.", PT_H2, True, INK)]],
                 sub="이 한 줄이 1시간의 중심축입니다.")
    draw_chrome(s, 3, "오프닝", "THE THESIS")

    # S04 발전사 4단계 — T3 타임라인
    s = blank_slide(prs)
    T3_timeline(s, "발전사 4단계 — 지금은 3단계",
                stages=[("1", "RULES", ["엑셀 매크로·VLOOKUP", "사람이 짠 규칙대로"]),
                        ("2", "ML", ["통계·머신러닝", "데이터에서 패턴"]),
                        ("3", "LLM", ["ChatGPT·Claude", "말 이어가는 두뇌 (지금)"]),
                        ("4", "AGENT", ["에이전트·AX", "도구 달아 맡긴다 (다음)"])],
                active=2, dashed=3,
                caption="에이전트는 새 두뇌가 아니라 3단계 위의 응용 단계  ·  트랜스포머·파라미터 언급 금지")
    draw_chrome(s, 4, "프롤로그", "EVOLUTION · 4 STAGES")

    # S05 ★ 개념지도 (완전 구현)
    s = blank_slide(prs)
    T4_concept_map(s)
    draw_chrome(s, 5, "프롤로그·큰 그림", "PROLOGUE · BIG PICTURE")

    # S06 가치 먼저 — T5 3카드
    s = blank_slide(prs)
    T5_cards(s, "AI가 잘 하는 일",
             cards=[("①", "분개장 요약", ["긴 분개장을", "3줄 핵심로 압축"]),
                    ("②", "긴 메일·회의록", ["핵심만", "3줄로 압축"]),
                    ("③", "조서 초안 골격", ["빈 템플릿을", "항목별로 채운 초안"])],
             caption="언어를 다루는 일 — 정말 빠르고 쓸모 있습니다")
    draw_chrome(s, 6, "M1 원리", "M1 · PRINCIPLE")

    # S07 환각 캡처 A — 의도적 왜곡 예시 표기
    s = blank_slide(prs)
    T6_capture(s, png_path=os.path.join("captures", "02_환각틀림.png"),
               h1="환각(캡처 A) — 조용한 환각",
               caption_lines=[[S("실존 조문(조세특례제한법 §7)인데", PT_CAPTION, False, BODY)],
                              [S("감면율·조문번호 미세 왜곡", PT_CAPTION, True, WARN)],
                              [S("얼핏 맞아 보여 더 위험", PT_CAPTION, False, BODY)],
                              [S(" ", PT_CAPTION)],
                              [S("뒤 30분에 Claude Code로 라이브 재현·정정", PT_CAPTION, False, ACCENT)],
                              [S("확률적으로 찍는 것이 아님 — 가장 자연스러운 답을 고를 뿐", PT_CAPTION, False, MUTE)]],
               note="※ 의도적으로 유발한 환각 예시(강사 제작 캡처). 환각은 재현이 안 돼 미리 캡처.",
               warn=True)
    draw_chrome(s, 7, "M1 원리", "M1 · HALLUCINATION A")

    # S08 환각 캡처 B — 결정타
    s = blank_slide(prs)
    T6_capture(s, png_path=os.path.join("captures", "03_정정답.png"),
               h1="환각(캡처 B) — 결정타: 감면누락",
               caption_lines=[[S("§7 감면을 ", PT_CAPTION, False, BODY)],
                              [S("“2025.12.31. 종료·폐지”", PT_CAPTION, True, WARN)],
                              [S("로 왜곡 (실제 2028.12.31. 연장)", PT_CAPTION, False, BODY)],
                              [S("납세자 감면 안내에 실렸다면?", PT_CAPTION, True, INK)],
                              [S("→ 받을 혜택 누락", PT_CAPTION, True, WARN)]],
               note="※ 의도적 왜곡 예시. 조문 수치는 law.go.kr 사전 검증 필수(정직성 세이프가드).",
               warn=True)
    draw_chrome(s, 8, "M1 원리", "M1 · HALLUCINATION B")

    # S09 원리 + 정규화 — '확률적으로 찍는다 아님' 부표기
    s = blank_slide(prs)
    T7_split(s, "왜 틀리는지 — 원리",
             "결정적 계산기",
             [[S("늘 같은 답", PT_BODY, True, INK)], [S("1+1=2, 누적과세구간 합산", PT_BODY_SM, False, BODY)],
              [S("로컬 코드 — 밖으로 안 나감", PT_CAPTION, False, MUTE)]],
             "AI",
             [[S("‘가장 자연스러운’ 답을 고른다", PT_BODY, True, INK)],
              [S("→ 정답과 다를 때가 환각", PT_BODY_SM, False, WARN)],
              [S("‘확률적으로 찍는다’가 아님 (금지 표현)", PT_CAPTION, False, MUTE)]],
             bottom="데이터를 정규화해 결정적 분석 → 그게 뒤 1시간의 뼈대")
    draw_chrome(s, 9, "M1 원리", "M1 · PRINCIPLE")

    # S10 프롬프트 — T5 칩 (4+3)
    s = blank_slide(prs)
    T5_chips(s, "물어보는 법이 결과를 바꾼다",
             top_chips=["역할", "맥락", "지시", "형식"],
             bottom_chips=["예시(few-shot)", "단계별로 생각해", "모르면 모른다고 해"],
             caption="예시 하나만 붙여도 결과가 다릅니다 — 뒤 30분에 직접 보여드립니다")
    draw_chrome(s, 10, "M1 원리", "M1 · PROMPT")

    # S11 UX→AX — T7 대비 + 화살표
    s = blank_slide(prs)
    T7_split(s, "UX → AX — 챗봇을 넘어 보조 인력",
             "UX",
             [[S("내가 매 단계 누른다", PT_BODY, True, INK)],
              [S("질문 → 답, 또 질문 → 답", PT_BODY_SM, False, BODY)],
              [S("챗봇 / 계산기", PT_CAPTION, False, MUTE)]],
             "AX",
             [[S("‘이 일 끝까지 가져와 줘’", PT_BODY, True, INK)],
              [S("AI가 찾고·모으고·정리해 초안까지", PT_BODY_SM, False, BODY)],
              [S("일을 끝까지 맡기는 보조 인력", PT_CAPTION, False, MUTE)]],
             bottom="챗봇이 아니라, 일을 끝까지 가져와 주는 보조 인력을 두는 단계",
             arrow=True)
    draw_chrome(s, 11, "M3 AX", "M3 · UX → AX")

    # S12 하네스 — T8 다이어그램
    s = blank_slide(prs)
    T8_harness(s, "하네스 — 실행 틀 + 고삐 2개",
               outer_label="HARNESS — 에이전트가 돌아가는 실행 틀 (도구·가드레일 포함)",
               inner_boxes=[("CLAUDE.md", "업무지시서 — 역할·출력 순서·말투·‘근거 없으면 인용 마’, 파일 하나로 내 기준 고정"),
                            ("Skill", "절차 매뉴얼 — 반복 업무를 명령 한 줄로 재사용 (예: /세법질의 제7조)")],
               side_note="이런 것도 있다 — 서브에이전트(역할 분담) · MCP(더존·홈택스 연결)",
               bottom="이 두 고삐를 뒤 30분에 제가 직접 만들어 보여드립니다")
    draw_chrome(s, 12, "M3 AX", "M3 · HARNESS")

    # S13 회계 분업선 — 기계(회색) vs 회계사(오렌지) 명암 대비
    s = blank_slide(prs)
    T7_split(s, "회계 분업선 — 기계 자리 vs 회계사 자리",
             "AX가 가져갈 자리",
             [[S("수임처별 대사(tie-out)", PT_BODY_SM, False, BODY)],
              [S("월결산 마감 진척 수집", PT_BODY_SM, False, BODY)],
              [S("조회서 PDF 처리·집계", PT_BODY_SM, False, BODY)],
              [S("조서 템플릿 채움", PT_BODY_SM, False, BODY)],
              [S("규칙은 있되 매번 손으로 엮어야 하는 일", PT_CAPTION, False, MUTE)]],
             "회계사가 가져갈 자리",
             [[S("판단", PT_BODY, True, INK)], [S("의심", PT_BODY, True, INK)],
              [S("결론", PT_BODY, True, INK)], [S("서명", PT_BODY, True, ACCENT)],
              [S("→ 이건 옮기지 않습니다", PT_CAPTION, True, INK)]],
             bottom="최종 책임은 서명 회계사 — 회계사 자리는 사라지지 않는다")
    draw_chrome(s, 13, "M3 AX", "M3 · DIVISION OF LABOR")

    # S14 결과물 — 감사 본류 미구축 세이프가드
    s = blank_slide(prs)
    T6_capture(s, png_path=os.path.join("captures", "07_결과물정지화면.png"),
               h1="이 고삐를 세팅하면",
               caption_lines=[[S("세법 질의 Skill로 포맷팅된 결과(더미)", PT_CAPTION, False, BODY)],
                              [S("반복 업무가 명령 한 줄로", PT_CAPTION, False, BODY)],
                              [S(" ", PT_CAPTION)],
                              [S("※ 감사 본류(위험평가·조서)는 미구축 — 본 캡처는 원천세/세법질의 영역", PT_CAPTION, True, WARN)]],
               note="직접 만드는 건 뒤 30분에 보여드립니다.")
    draw_chrome(s, 14, "M3 AX", "M3 · OUTPUT")

    # S15 보안 — T8 보안형 (뼈대 + 별도검토 3종)
    s = blank_slide(prs)
    T8_security(s, "보안 — ‘접어두고’가 아니라 ‘푼 방식’",
                base_lines=[[S("한 줄 뼈대:", PT_BODY, True, ACCENT)],
                            [S("숫자·계산은 전부 로컬 코드(밖으로 안 나감) · AI엔 비식별 최소 정보 · 더 민감하면 사내 로컬 모델 + 무보존(무학습·무저장) 계약", PT_BODY_SM, False, BODY)]],
                risks=[("①", "재식별 위험 — 사업자번호+금액+귀속월이면 식별 가능"),
                       ("②", "미공시정보 — 상장·코스닥 수감회사 결산 자료는 로컬 강제"),
                       ("③", "감리·증적 보존 — AI 산출물도 감사증적으로 보존")],
                side_note="다시, 말로만 숫자·사실은 틀리니까 → 숫자는 코드, 사실은 검색에 (환각 방지)")
    draw_chrome(s, 15, "M4 신뢰·책임", "M4 · SECURITY")

    # S16 기계 vs 사람 + 책임소재
    s = blank_slide(prs)
    T7_split(s, "기계 vs 사람 + 책임소재",
             "기계",
             [[S("취합", PT_BODY, True, BODY)],
              [S("정리", PT_BODY, True, BODY)],
              [S("초안", PT_BODY, True, BODY)]],
             "사람",
             [[S("판단", PT_BODY, True, INK)],
              [S("의심 · 결론", PT_BODY, True, INK)],
              [S("서명", PT_BODY, True, ACCENT)]],
             bottom="AI 산출물의 최종 책임은 서명한 회계사 — 회계사의 자리가 사라지지 않습니다")
    draw_chrome(s, 16, "M4 신뢰·책임", "M4 · MACHINE vs HUMAN")
    add_box(s, 6, 90.5, 88, 3, [[S("쓰는 동료와 안 쓰는 동료의 격차가 1~2년 안에 벌어집니다", PT_CAPTION, False, MUTE)]],
            align=PP_ALIGN.CENTER)

    # S17 후반부 진입표 (pilot 티저)
    s = blank_slide(prs)
    T9_table(s, "개념은 여기까지, 이제 Claude Code를 켭니다",
             ["비트", "시간", "무엇"],
             [["비트① 환각→정정", "7분", "캡처 환각을 Claude Code 안에서 라이브 재현, 4요소+3강화로 정정"],
              ["★ 비트② CLAUDE.md", "10분", "빈 파일에 업무지시서 한 장 → 내 기준·말투로 답 + 폰 따라하기 (결정적 한 장면)"],
              ["비트③ Skill 실행", "4분", "/세법질의 제7조 한 줄 → 4단 포맷 (만드는 법은 유료)"]],
             highlight_row=1,
             bottom="그리고 뒤이어 1시간 — 완성된 회계 감사 시스템(위험평가·조서·정산표)을 돌려보여드립니다")
    draw_chrome(s, 17, "후반부 진입", "ENTER · LIVE DEMO")

    prs.save(out)
    print(f"✓ saved: {out}  ({len(prs.slides._sldIdLst)} slides)")

if __name__ == "__main__":
    build()
